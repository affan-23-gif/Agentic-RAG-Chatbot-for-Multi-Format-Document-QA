# agents/ingestion_agent.py
import os
from pypdf import PdfReader
from docx import Document
import pandas as pd
from pptx import Presentation
from mcp import MCPMessage # Assuming mcp.py is in the parent directory or correctly imported

class IngestionAgent:
    def __init__(self, coordinator_queue=None):
        self.coordinator_queue = coordinator_queue # For sending messages back to coordinator

    def parse_pdf(self, file_path):
        text = ""
        with open(file_path, "rb") as file:
            reader = PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() or ""
        return text

    def parse_docx(self, file_path):
        doc = Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])

    def parse_pptx(self, file_path):
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def parse_csv(self, file_path):
        df = pd.read_csv(file_path)
        return df.to_string() # Convert DataFrame to string for indexing

    def parse_txt_md(self, file_path):
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read()

    def process_document(self, file_path: str, trace_id: str):
        file_extension = os.path.splitext(file_path)[1].lower()
        content = ""
        document_type = ""

        try:
            if file_extension == ".pdf":
                content = self.parse_pdf(file_path)
                document_type = "pdf"
            elif file_extension == ".docx":
                content = self.parse_docx(file_path)
                document_type = "docx"
            elif file_extension == ".pptx":
                content = self.parse_pptx(file_path)
                document_type = "pptx"
            elif file_extension == ".csv":
                content = self.parse_csv(file_path)
                document_type = "csv"
            elif file_extension in [".txt", ".md"]:
                content = self.parse_txt_md(file_path)
                document_type = "txt_md"
            else:
                print(f"Unsupported file type: {file_extension}")
                return None

            # Basic chunking (can be made more sophisticated)
            chunks = self.chunk_text(content, chunk_size=1000, overlap=100)
            print(f"Parsed {document_type} and created {len(chunks)} chunks.")

            # Send parsed content to RetrievalAgent via MCP
            if self.coordinator_queue:
                msg = MCPMessage(
                    sender="IngestionAgent",
                    receiver="RetrievalAgent",
                    msg_type="DOCUMENT_PARSED",
                    trace_id=trace_id,
                    payload={
                        "document_name": os.path.basename(file_path),
                        "document_type": document_type,
                        "chunks": chunks
                    }
                )
                self.coordinator_queue.put(msg)
            return chunks

        except Exception as e:
            print(f"Error processing {file_path}: {e}")
            return None

    def chunk_text(self, text: str, chunk_size: int, overlap: int):
        # Simple character-based chunking. For production, consider token-based chunking
        # and more advanced strategies (e.g., Langchain's RecursiveCharacterTextSplitter).
        chunks = []
        if not text:
            return chunks
        for i in range(0, len(text), chunk_size - overlap):
            chunks.append(text[i:i + chunk_size])
        return chunks